"""
Voxaris x Family Orlando Dentistry — Pitch Deck Generator
Generates a polished 16-slide + sources slide .pptx
using python-pptx with enhanced visual design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ─── Brand palette ───────────────────────────────────────────
BLUE   = RGBColor(22, 120, 212)
NAVY   = RGBColor(15, 39, 64)
WHITE  = RGBColor(255, 255, 255)
SLATE  = RGBColor(90, 108, 128)
TEAL   = RGBColor(14, 165, 164)
BG     = RGBColor(245, 249, 253)
LIGHT_BLUE = RGBColor(220, 235, 250)
ACCENT = RGBColor(59, 190, 160)   # secondary green accent

# ─── Slide dimensions (16:9 widescreen) ─────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Speaker notes for each slide ────────────────────────────
SPEAKER_NOTES = [
    "Open by framing this as a conversion and workload solution, not a generic AI tool. The practice already offers strong clinical services; the missing piece is always-on access.",
    "This is a classic high-trust, family-run office with real production upside if access gets modernized.",
    "This is not a low-value hygiene-only office. Their online demand should be treated like premium revenue opportunities.",
    "The site may be adequate as an information page, but it is underbuilt as a conversion machine.",
    "This is the bridge between their current operating model and why Voxaris is timely.",
    "This validates the category while opening space for Voxaris differentiation.",
    "Position Voxaris as higher-touch, higher-conversion, and more presentation-ready.",
    "Show how the flow handles both premium elective and urgent service types.",
    "Make this the money slide. Most private practices underuse their own database.",
    "This should feel credible, not inflated. Explain that database quality and segmentation matter.",
    "Tie the transformation directly to the owner's experience of missed after-hours demand.",
    "Owners want proof, not abstraction. The dashboard is how Voxaris stays accountable.",
    "Keep these as illustrative projections, not guarantees.",
    "This slide matters because this is a family-run office, not a faceless chain.",
    "Show a low-friction adoption path. Start with the fastest pain relief.",
    "End with momentum and specificity.",
]


# ═══════════════════════════════════════════════════════════════
#  SLIDE DATA
# ═══════════════════════════════════════════════════════════════
slides_data = [
    {
        "num": 1,
        "title": "Family Orlando Dentistry\nx Voxaris",
        "subtitle": "A 24/7 digital front desk for a phone-first private practice",
        "bullets": [
            "Custom demo concept for Family Orlando Dentistry",
            "Focus: inbound capture, reactivation, and V\u00b7FACE-led trust",
            "Prepared for Voxaris sales presentation",
        ],
        "layout": "title",
    },
    {
        "num": 2,
        "title": "What makes this practice special",
        "bullets": [
            "Husband-and-wife doctor team with 55+ combined years of experience",
            "Ocoee private practice with limited weekday access hours",
            "General, restorative, and cosmetic positioning already established",
        ],
        "accent_label": "FAMILY ORLANDO",
    },
    {
        "num": 3,
        "title": "High-value services are already on the table",
        "bullets": [
            "Invisalign with free consultation",
            "Same-day / single-visit crowns, implants, whitening, emergency care",
            "Root canals, extractions, bone grafting, oral appliances, Botox / TMJ support",
        ],
        "accent_label": "FAMILY ORLANDO",
    },
    {
        "num": 4,
        "title": "The website acts like a brochure,\nnot a front desk",
        "bullets": [
            "Homepage pushes phone-first behavior",
            "No contact form and no booking form on the contact page",
            "Auto-draft page slugs weaken trust and polish",
        ],
        "accent_label": "FAMILY ORLANDO",
    },
    {
        "num": 5,
        "title": "Patients still book by phone,\nbut they expect modern access",
        "bullets": [
            "Dental booking is still heavily phone-driven",
            "A large share of booking activity happens outside business hours when digital access exists",
            "Limited office availability creates a conversion gap",
        ],
        "layout": "stats",
        "stats": [
            {"value": "85%", "label": "phone-booked\nappointments"},
            {"value": "47%", "label": "booked outside\nbusiness hours"},
            {"value": "4 days", "label": "open per week\n(no Wed/Sat/Sun)"},
        ],
    },
    {
        "num": 6,
        "title": "The dental AI market has\nvalidated the pain",
        "bullets": [
            "Arini, Newton, Viva, and Yobi all market 24/7 answering and recall/reactivation",
            "The category proves strong demand for inbound and outbound automation",
            "Most still do not lead with photorealistic conversational video",
        ],
    },
    {
        "num": 7,
        "title": "Where Voxaris wins",
        "bullets": [
            "V\u00b7TEAMS: receptionist \u2192 qualifier \u2192 specialist \u2192 closer",
            "V\u00b7OUTBOUND: database revival and follow-up",
            "V\u00b7FACE: photorealistic digital concierge layer",
        ],
        "layout": "three_cards",
        "cards": [
            {"icon": "\u260e", "heading": "V\u00b7TEAMS", "body": "Multi-stage routing\nfor every call type"},
            {"icon": "\U0001f4e4", "heading": "V\u00b7OUTBOUND", "body": "Turn dormant lists\ninto booked visits"},
            {"icon": "\U0001f3ac", "heading": "V\u00b7FACE", "body": "Photorealistic avatar\nfront-door presence"},
        ],
    },
    {
        "num": 8,
        "title": "New patient and emergency capture",
        "bullets": [
            "Instant handling for Invisalign, whitening, implants, and emergencies",
            "No repeated context for the patient",
            "Priority routing for urgent and high-intent requests",
        ],
    },
    {
        "num": 9,
        "title": "The hidden revenue is\nin the database",
        "bullets": [
            "Overdue hygiene reactivation",
            "Warm Invisalign follow-up",
            "Unscheduled treatment and cancellation gap recovery",
        ],
    },
    {
        "num": 10,
        "title": "What good reactivation\ncan look like",
        "bullets": [
            "10\u201315% overdue-list reactivation per campaign is a solid benchmark",
            "15\u201325% recovery is a strong upside case in the right setup",
            "Direct outreach materially improves return rates",
        ],
        "layout": "benchmark",
        "benchmarks": [
            {"label": "Conservative", "low": 10, "high": 15},
            {"label": "Strong setup", "low": 15, "high": 25},
        ],
    },
    {
        "num": 11,
        "title": "From static website to\nliving front desk",
        "bullets": [
            "V\u00b7FACE as the front-door concierge",
            "Service qualification for Invisalign, emergencies, same-day crowns, whitening, and implants",
            "Always-on capture even when the office is closed",
        ],
    },
    {
        "num": 12,
        "title": "What the owner gets to see",
        "bullets": [
            "Calls answered",
            "Appointments booked",
            "Missed-call recovery",
            "Revenue lift and service-line pipeline",
        ],
        "layout": "four_kpis",
        "kpis": [
            {"icon": "\U0001f4de", "label": "Calls\nAnswered"},
            {"icon": "\U0001f4c5", "label": "Appointments\nBooked"},
            {"icon": "\U0001f504", "label": "Missed-Call\nRecovery"},
            {"icon": "\U0001f4c8", "label": "Revenue\nPipeline"},
        ],
    },
    {
        "num": 13,
        "title": "A simple 90-day outcome model",
        "bullets": [
            "Recover after-hours and overflow demand",
            "Increase high-LTV consult requests",
            "Reactivate part of the dormant database",
            "Reduce staff and doctor phone burden",
        ],
    },
    {
        "num": 14,
        "title": "This is not just about revenue",
        "bullets": [
            "Fewer personal-cell interruptions",
            "Less voicemail pressure",
            "Better patient experience",
            "More protected family time",
        ],
    },
    {
        "num": 15,
        "title": "How this could roll out",
        "bullets": [
            "Phase 1 \u2014 Inbound coverage and emergency handling",
            "Phase 2 \u2014 V\u00b7FACE website layer and high-value service qualification",
            "Phase 3 \u2014 Outbound reactivation and dashboard optimization",
        ],
        "layout": "phases",
    },
    {
        "num": 16,
        "title": "The next step",
        "bullets": [
            "Build the Family Orlando Dentistry demo environment",
            "Show live inbound, outbound, V\u00b7FACE, and dashboard experiences",
            "Align on launch scope and first KPI targets",
        ],
        "layout": "cta",
    },
]

# ─── Source URLs ──────────────────────────────────────────────
SOURCES = [
    "https://familyorlandodentistry.com/",
    "https://familyorlandodentistry.com/auto-draft-5/",
    "https://familyorlandodentistry.com/auto-draft-3/",
    "https://familyorlandodentistry.com/auto-draft-4/",
    "https://familyorlandodentistry.com/auto-draft-2/",
    "https://www.arini.ai",
    "https://www.arini.ai/blog/dental-call-centers-vs-ai-receptionists",
    "https://joinnewton.com/product/ai-receptionist",
    "https://www.getviva.ai",
    "https://dental.yobi.com",
    "https://pmc.ncbi.nlm.nih.gov/articles/PMC8796867/",
    "https://blog.softwareofexcellence.com/en-gb/why-online-booking-is-the-new-standard-in-dentistry",
    "https://mvpmailhouse.com/resources/blog/how-to-increase-patient-retention-in-dentistry",
    "https://mybcat.com/blog/reactivation-campaigns-dormant-patients-2026/",
    "https://practiceanalytics.com/the-impact-of-hygiene-reactivation/",
]


# ═══════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_navy_band(slide, height=Inches(0.55)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    return shape


def add_bottom_accent(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def add_slide_number(slide, num, total=16):
    tx = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45), Inches(1.0), Inches(0.3)
    )
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = SLATE


def add_title_text(slide, text, left=Inches(0.8), top=Inches(0.75),
                   width=Inches(8.0), height=Inches(1.2),
                   size=Pt(30), color=NAVY, bold=True):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Aptos Display"
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(2)


def add_subtitle_text(slide, text, left=Inches(0.82), top=Inches(1.85),
                      width=Inches(7.8)):
    tx = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.color.rgb = SLATE


def add_bullet_list(slide, bullets, left=Inches(1.0), top=Inches(2.5),
                    width=Inches(7.2), height=Inches(4.2),
                    size=Pt(20), color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2022  {bullet}"
        run.font.name = "Aptos"
        run.font.size = size
        run.font.color.rgb = color
        p.space_after = Pt(14)
        p.line_spacing = Pt(28)


def add_right_panel(slide, accent_text="VOXARIS"):
    # Outer panel
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.0), Inches(1.2), Inches(3.8), Inches(5.6)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT_BLUE
    panel.line.color.rgb = BLUE
    panel.line.width = Pt(1)

    # Top label
    top = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.3), Inches(1.55), Inches(3.2), Inches(1.0)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    tf = top.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = accent_text
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Feature list
    body = slide.shapes.add_textbox(Inches(9.4), Inches(2.85), Inches(3.0), Inches(3.4))
    tf2 = body.text_frame
    tf2.word_wrap = True
    items = [
        "24/7 coverage",
        "V\u00b7TEAMS routing",
        "V\u00b7OUTBOUND reactivation",
        "V\u00b7FACE presence",
    ]
    for idx, item in enumerate(items):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        run = p.add_run()
        run.text = f"\u2713  {item}"
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.color.rgb = NAVY
        p.space_after = Pt(14)


def add_speaker_notes(slide, notes_text):
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
    else:
        notes_slide = slide.notes_slide  # creates it
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_divider(slide, top):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), top, Inches(2.5), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


# ═══════════════════════════════════════════════════════════════
#  BUILD THE DECK
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

for idx, data in enumerate(slides_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    layout = data.get("layout", "standard")
    num = data["num"]

    # ─── Title slide ──────────────────────────────────────────
    if layout == "title":
        set_bg(slide, NAVY)
        # Big white title
        add_title_text(
            slide, data["title"],
            left=Inches(0.9), top=Inches(1.5),
            width=Inches(7.5), height=Inches(2.2),
            size=Pt(42), color=WHITE, bold=True,
        )
        # Subtitle
        tx = slide.shapes.add_textbox(Inches(0.92), Inches(3.6), Inches(7.5), Inches(0.6))
        p = tx.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = data["subtitle"]
        run.font.name = "Aptos"
        run.font.size = Pt(18)
        run.font.color.rgb = LIGHT_BLUE

        # Right accent panel
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.2), Inches(1.0), Inches(3.5), Inches(5.2)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = BLUE
        panel.line.fill.background()
        # Bullets inside panel
        bx = slide.shapes.add_textbox(Inches(9.5), Inches(1.6), Inches(3.0), Inches(4.0))
        tf = bx.text_frame
        tf.word_wrap = True
        for bi, b in enumerate(data["bullets"]):
            p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"\u2022  {b}"
            run.font.name = "Aptos"
            run.font.size = Pt(15)
            run.font.color.rgb = WHITE
            p.space_after = Pt(14)

        # Bottom bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()

    # ─── Stats layout (slide 5) ──────────────────────────────
    elif layout == "stats":
        set_bg(slide)
        add_navy_band(slide)
        add_title_text(slide, data["title"])
        add_divider(slide, Inches(1.55))

        stats = data.get("stats", [])
        card_w = Inches(3.4)
        gap = Inches(0.4)
        total_w = card_w * len(stats) + gap * (len(stats) - 1)
        start_x = (SLIDE_W - total_w) / 2
        for si, stat in enumerate(stats):
            x = start_x + si * (card_w + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), card_w, Inches(3.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = LIGHT_BLUE
            card.line.width = Pt(1)
            card.shadow.inherit = False

            # Big number
            vtx = slide.shapes.add_textbox(x, Inches(2.9), card_w, Inches(1.2))
            vp = vtx.text_frame.paragraphs[0]
            vp.alignment = PP_ALIGN.CENTER
            vr = vp.add_run()
            vr.text = stat["value"]
            vr.font.name = "Aptos Display"
            vr.font.size = Pt(44)
            vr.font.bold = True
            vr.font.color.rgb = BLUE

            # Label
            ltx = slide.shapes.add_textbox(x + Inches(0.3), Inches(4.1), card_w - Inches(0.6), Inches(1.2))
            ltf = ltx.text_frame
            ltf.word_wrap = True
            for li, line in enumerate(stat["label"].split("\n")):
                lp = ltf.paragraphs[0] if li == 0 else ltf.add_paragraph()
                lp.alignment = PP_ALIGN.CENTER
                lr = lp.add_run()
                lr.text = line
                lr.font.name = "Aptos"
                lr.font.size = Pt(16)
                lr.font.color.rgb = NAVY

        add_bottom_accent(slide)
        add_slide_number(slide, num)

    # ─── Three cards layout (slide 7 - Voxaris wins) ─────────
    elif layout == "three_cards":
        set_bg(slide)
        add_navy_band(slide)
        add_title_text(slide, data["title"])
        add_divider(slide, Inches(1.55))

        cards = data.get("cards", [])
        card_w = Inches(3.5)
        gap = Inches(0.45)
        total_w = card_w * len(cards) + gap * (len(cards) - 1)
        start_x = (SLIDE_W - total_w) / 2
        for ci, card_data in enumerate(cards):
            x = start_x + ci * (card_w + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), card_w, Inches(4.0)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = LIGHT_BLUE
            card.line.width = Pt(1)

            # Icon circle
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - Inches(1.0)) / 2, Inches(2.8),
                Inches(1.0), Inches(1.0)
            )
            circ.fill.solid()
            circ.fill.fore_color.rgb = NAVY
            circ.line.fill.background()
            tf_c = circ.text_frame
            tf_c.clear()
            p_c = tf_c.paragraphs[0]
            p_c.alignment = PP_ALIGN.CENTER
            r_c = p_c.add_run()
            r_c.text = card_data["icon"]
            r_c.font.size = Pt(28)
            r_c.font.color.rgb = WHITE

            # Heading
            htx = slide.shapes.add_textbox(x, Inches(4.0), card_w, Inches(0.5))
            hp = htx.text_frame.paragraphs[0]
            hp.alignment = PP_ALIGN.CENTER
            hr = hp.add_run()
            hr.text = card_data["heading"]
            hr.font.name = "Aptos Display"
            hr.font.size = Pt(20)
            hr.font.bold = True
            hr.font.color.rgb = NAVY

            # Body
            btx = slide.shapes.add_textbox(x + Inches(0.3), Inches(4.6), card_w - Inches(0.6), Inches(1.5))
            btf = btx.text_frame
            btf.word_wrap = True
            for bli, bline in enumerate(card_data["body"].split("\n")):
                bp = btf.paragraphs[0] if bli == 0 else btf.add_paragraph()
                bp.alignment = PP_ALIGN.CENTER
                br = bp.add_run()
                br.text = bline
                br.font.name = "Aptos"
                br.font.size = Pt(15)
                br.font.color.rgb = SLATE

        add_bottom_accent(slide)
        add_slide_number(slide, num)

    # ─── Benchmark layout (slide 10) ─────────────────────────
    elif layout == "benchmark":
        set_bg(slide)
        add_navy_band(slide)
        add_title_text(slide, data["title"])
        add_divider(slide, Inches(1.55))
        add_bullet_list(slide, data["bullets"], top=Inches(2.2), height=Inches(2.5))

        benchmarks = data.get("benchmarks", [])
        bar_left = Inches(1.0)
        bar_top_start = Inches(5.0)
        max_w = Inches(10.0)
        for bi, bm in enumerate(benchmarks):
            y = bar_top_start + bi * Inches(0.85)
            # Label
            ltx = slide.shapes.add_textbox(bar_left, y - Inches(0.3), Inches(2.0), Inches(0.3))
            lp = ltx.text_frame.paragraphs[0]
            lr = lp.add_run()
            lr.text = bm["label"]
            lr.font.name = "Aptos"
            lr.font.size = Pt(14)
            lr.font.color.rgb = NAVY
            lr.font.bold = True

            # Background bar
            bg_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left, y, max_w, Inches(0.35)
            )
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = LIGHT_BLUE
            bg_bar.line.fill.background()

            # Fill bar
            fill_w = int(max_w * bm["high"] / 30)  # scale to 30% max
            fill_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left, y, fill_w, Inches(0.35)
            )
            fill_bar.fill.solid()
            fill_bar.fill.fore_color.rgb = BLUE if bi == 0 else TEAL
            fill_bar.line.fill.background()

            # Percentage label
            ptx = slide.shapes.add_textbox(bar_left + fill_w + Inches(0.15), y, Inches(1.5), Inches(0.35))
            ptf = ptx.text_frame
            ptf.paragraphs[0].alignment = PP_ALIGN.LEFT
            pp = ptf.paragraphs[0]
            pr = pp.add_run()
            pr.text = f"{bm['low']}\u2013{bm['high']}%"
            pr.font.name = "Aptos Display"
            pr.font.size = Pt(16)
            pr.font.bold = True
            pr.font.color.rgb = NAVY

        add_bottom_accent(slide)
        add_slide_number(slide, num)

    # ─── Four KPI cards (slide 12) ───────────────────────────
    elif layout == "four_kpis":
        set_bg(slide)
        add_navy_band(slide)
        add_title_text(slide, data["title"])
        add_divider(slide, Inches(1.55))

        kpis = data.get("kpis", [])
        card_w = Inches(2.6)
        gap = Inches(0.35)
        total_w = card_w * len(kpis) + gap * (len(kpis) - 1)
        start_x = (SLIDE_W - total_w) / 2
        for ki, kpi in enumerate(kpis):
            x = start_x + ki * (card_w + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), card_w, Inches(4.0)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = LIGHT_BLUE
            card.line.width = Pt(1)

            # Icon
            itx = slide.shapes.add_textbox(x, Inches(3.0), card_w, Inches(1.0))
            ip = itx.text_frame.paragraphs[0]
            ip.alignment = PP_ALIGN.CENTER
            ir = ip.add_run()
            ir.text = kpi["icon"]
            ir.font.size = Pt(40)

            # Label
            ktx = slide.shapes.add_textbox(x + Inches(0.2), Inches(4.2), card_w - Inches(0.4), Inches(1.5))
            ktf = ktx.text_frame
            ktf.word_wrap = True
            for kli, kline in enumerate(kpi["label"].split("\n")):
                kp = ktf.paragraphs[0] if kli == 0 else ktf.add_paragraph()
                kp.alignment = PP_ALIGN.CENTER
                kr = kp.add_run()
                kr.text = kline
                kr.font.name = "Aptos Display"
                kr.font.size = Pt(18)
                kr.font.bold = True
                kr.font.color.rgb = NAVY

        add_bottom_accent(slide)
        add_slide_number(slide, num)

    # ─── Phases timeline (slide 15) ──────────────────────────
    elif layout == "phases":
        set_bg(slide)
        add_navy_band(slide)
        add_title_text(slide, data["title"])
        add_divider(slide, Inches(1.55))

        phases = data["bullets"]
        card_w = Inches(3.5)
        gap = Inches(0.45)
        total_w = card_w * 3 + gap * 2
        start_x = (SLIDE_W - total_w) / 2
        colors = [BLUE, TEAL, NAVY]
        for pi, phase in enumerate(phases):
            x = start_x + pi * (card_w + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.8), card_w, Inches(3.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = colors[pi]
            card.line.fill.background()

            # Phase number circle
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - Inches(0.8)) / 2, Inches(3.1),
                Inches(0.8), Inches(0.8)
            )
            circ.fill.solid()
            circ.fill.fore_color.rgb = WHITE
            circ.line.fill.background()
            tf_c = circ.text_frame
            tf_c.clear()
            p_c = tf_c.paragraphs[0]
            p_c.alignment = PP_ALIGN.CENTER
            r_c = p_c.add_run()
            r_c.text = str(pi + 1)
            r_c.font.name = "Aptos Display"
            r_c.font.size = Pt(24)
            r_c.font.bold = True
            r_c.font.color.rgb = colors[pi]

            # Phase text
            ptx = slide.shapes.add_textbox(x + Inches(0.3), Inches(4.2), card_w - Inches(0.6), Inches(1.8))
            ptf = ptx.text_frame
            ptf.word_wrap = True
            # Remove "Phase N — " prefix for card display
            phase_text = phase.split(" \u2014 ", 1)[-1] if " \u2014 " in phase else phase
            pp = ptf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            pr = pp.add_run()
            pr.text = phase_text
            pr.font.name = "Aptos"
            pr.font.size = Pt(16)
            pr.font.color.rgb = WHITE

        # Connecting arrows
        for ai in range(2):
            x_arrow = start_x + (ai + 1) * (card_w + gap) - gap / 2 - Inches(0.15)
            atx = slide.shapes.add_textbox(x_arrow, Inches(4.2), Inches(0.3), Inches(0.5))
            ap = atx.text_frame.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = "\u25b6"
            ar.font.size = Pt(18)
            ar.font.color.rgb = SLATE

        add_bottom_accent(slide)
        add_slide_number(slide, num)

    # ─── CTA slide (slide 16) ────────────────────────────────
    elif layout == "cta":
        set_bg(slide, NAVY)
        add_title_text(
            slide, data["title"],
            left=Inches(0.9), top=Inches(1.8),
            width=Inches(11.0), height=Inches(1.2),
            size=Pt(40), color=WHITE, bold=True,
        )
        # Divider
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.9), Inches(2.9), Inches(3.0), Pt(4)
        )
        div.fill.solid()
        div.fill.fore_color.rgb = TEAL
        div.line.fill.background()

        add_bullet_list(
            slide, data["bullets"],
            left=Inches(1.1), top=Inches(3.3),
            width=Inches(10.0), height=Inches(3.0),
            size=Pt(22), color=WHITE,
        )

        # CTA button shape
        btn = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.1), Inches(5.8), Inches(3.5), Inches(0.7)
        )
        btn.fill.solid()
        btn.fill.fore_color.rgb = BLUE
        btn.line.fill.background()
        tf = btn.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Book Your Live Demo"
        run.font.name = "Aptos Display"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Practice phone
        ptx = slide.shapes.add_textbox(Inches(5.0), Inches(5.85), Inches(4.0), Inches(0.5))
        pp = ptx.text_frame.paragraphs[0]
        pr = pp.add_run()
        pr.text = "(407) 877-9003  |  familyorlandodentistry.com"
        pr.font.name = "Aptos"
        pr.font.size = Pt(14)
        pr.font.color.rgb = LIGHT_BLUE

        # Bottom teal bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()

    # ─── Standard layout (default) ───────────────────────────
    else:
        set_bg(slide)
        add_navy_band(slide)
        add_title_text(slide, data["title"])
        add_divider(slide, Inches(1.55))
        add_bullet_list(slide, data["bullets"])
        accent = data.get("accent_label", "VOXARIS")
        add_right_panel(slide, accent)
        add_bottom_accent(slide)
        add_slide_number(slide, num)

    # Speaker notes
    if idx < len(SPEAKER_NOTES):
        add_speaker_notes(slide, SPEAKER_NOTES[idx])


# ─── Source slide (bonus) ────────────────────────────────────
src_slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(src_slide)
add_navy_band(src_slide)
add_title_text(src_slide, "Selected source URLs")
add_divider(src_slide, Inches(1.55))
add_subtitle_text(src_slide, "Include these in presenter notes or follow-up materials as needed")

src_box = src_slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.8))
src_tf = src_box.text_frame
src_tf.word_wrap = True
for i, src in enumerate(SOURCES):
    p = src_tf.paragraphs[0] if i == 0 else src_tf.add_paragraph()
    run = p.add_run()
    run.text = src
    run.font.name = "Aptos"
    run.font.size = Pt(13)
    run.font.color.rgb = BLUE
    run.font.underline = True
    p.space_after = Pt(5)

add_bottom_accent(src_slide)
add_slide_number(src_slide, 17, 17)
add_speaker_notes(src_slide, "These URLs were used during research. Verify before citing publicly.")


# ─── Save ────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "Voxaris_FamilyOrlandoDentistry_Demo_Deck.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
